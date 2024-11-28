import logging.config
import os
import time
import shutil
from datetime import datetime
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
import anthropic
from PIL import Image
import PyPDF2
import mimetypes
from dotenv import load_dotenv
from mutagen import File as MutagenFile
import base64
import mimetypes
# import textract
from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation
from openpyxl import load_workbook
import re
import nltk
import os
import logging

load_dotenv()

# Configure logging
logging.basicConfig(
    filename='anthropic_main.log',  # Save logs to a file
    filemode='a',  # Append mode
    format='%(asctime)s - %(levelname)s - %(message)s',  # Log format
    level=logging.INFO  # Log level
)

if not os.path.exists("/Users/joey/nltk_data/corpora/words") or not os.path.exists("/Users/joey/nltk_data/corpora/names"):
    # Download necessary NLTK data (run this once)
    nltk.download('words')
    nltk.download('names')

# Configure the directories to watch
WATCH_DIRECTORIES = [
    '/Users/joey/Downloads',
    '/Users/joey/Downloads/WebBrowser',
    '/Users/joey/Library/Mobile Documents/com~apple~CloudDocs/CleanShotProExports',
]

# File type mappings
FILE_TYPES = {
    'img': ['.jpg', '.jpeg', '.png', '.gif', '.bmp'],
    'pdf': ['.pdf'],
    'vid': ['.mp4', '.avi', '.mov', '.mkv'],
    'doc': ['.doc', '.docx', '.txt', '.rtf'],
    'ppt': ['.ppt', '.pptx'],
    'sheet': ['.xls', '.xlsx', '.csv']
}

MAX_CHARS = 30
MAX_TOKENS = 300


class FileHandler(FileSystemEventHandler):
    def on_created(self, event):
        self.anthropic_client = anthropic.Anthropic(
            api_key=os.getenv("ANHTROPIC_API_KEY"))
        if not event.is_directory:
            self.process_file(event.src_path)

    def process_file(self, file_path):
        file_ext = os.path.splitext(file_path)[1].lower()
        file_type = next((t for t, exts in FILE_TYPES.items()
                         if file_ext in exts), None)

        if file_type:
            summary_name = self.get_summary(file_path, file_type)
            new_name_dict = self.create_new_name(
                file_path, summary_name, file_type)
            new_full_path = new_name_dict["final_path"]
            new_file_name = new_name_dict["final_name"]
            old_filename = new_name_dict["old_filename"]
            directory = new_name_dict["directory"]
            self.move_and_rename_file(file_path, new_full_path, file_type)

    def create_new_name(self, original_file_path: str, summary_name: str, file_type: str):
        directory, old_filename = os.path.split(original_file_path)
        file_name, file_extension = os.path.splitext(old_filename)

        # Extract string content from summary_name if it's not already a string
        if not isinstance(summary_name, str):
            summary_name = str(summary_name)

        # Clean the summary_name to ensure it's filesystem-friendly
        summary_name = ''.join(
            c for c in summary_name if c.isalnum() or c in (' ', '-', '_')).strip()
        # Ensure it's no longer than MAX_CHARS characters
        summary_name = summary_name[:MAX_CHARS]

        date_stamp = datetime.now().strftime("%d-%m-%Y %H.%M.%S")
        if file_type == 'pdf':
            # For PDFs, use date-time stamp and summary name
            new_name = f"{date_stamp}_{summary_name}{file_extension}"
        elif file_type == 'img':
            # For images, append the summary to the existing name
            # new_name = f"{date_stamp}_{file_name}_{summary_name}{file_extension}"
            new_name = f"{date_stamp}_{summary_name}_{file_name}{file_extension}"
        else:
            # For other file types, use a combination of original name and summary
            new_name = f"{file_name[:(MAX_CHARS - 10)]}_{summary_name}{file_extension}"

        # Ensure the new name is unique in the directory
        counter = 1
        final_path = os.path.join(directory, new_name)
        while os.path.exists(final_path):
            if file_type == 'pdf':
                new_name = f"{date_stamp}_{summary_name}_{counter}{file_extension}"
            elif file_type == 'img':
                new_name = f"{file_name}_{summary_name}_{counter}{file_extension}"
            else:
                new_name = f"{file_name[:(MAX_CHARS - 10)]}_{summary_name}_{counter}{file_extension}"
            final_path = os.path.join(directory, new_name)
            counter += 1

        return {
            "final_path": final_path,
            "old_filename": old_filename,
            "directory": directory,
            "file_name": file_name,
            "file_extension": file_extension
        }

    def move_and_rename_file(self, original_path, new_path, file_type):
        try:
            # Create the subdirectory if it doesn't exist
            new_dir = os.path.join(os.path.dirname(original_path), file_type)
            os.makedirs(new_dir, exist_ok=True)

            # Adjust the new_path to be in the subdirectory
            new_filename = os.path.basename(new_path)
            final_path = os.path.join(new_dir, new_filename)

            # Move and rename the file
            shutil.move(original_path, final_path)
            logging.info(
                f"File moved and renamed: {original_path} -> {final_path}")
        except Exception as e:
            logging.info(
                f"Error moving and renaming file {original_path}: {e}")

    def get_summary(self, file_path, file_type):
        if file_type == 'img':
            content_block_list = self.summarize_image(file_path)
            return content_block_list[0].text
        elif file_type == 'pdf':
            content_block_list = self.summarize_pdf(file_path)
            return content_block_list[0].text
        elif file_type == 'vid':
            return self.get_video_metadata(file_path)
        else:
            return self.summarize_document(file_path)

    def summarize_image(self, image_path):
        mime_type, _ = mimetypes.guess_type(image_path)
        if not mime_type or not mime_type.startswith('image/'):
            mime_type = 'application/octet-stream'

        with open(image_path, "rb") as f:
            image_data = f.read()

        response = self.anthropic_client.messages.create(
            model="claude-3-sonnet-20240229",
            max_tokens=MAX_TOKENS,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": base64.b64encode(image_data).decode()
                            }
                        },
                        {
                            "type": "text",
                            "text": F"Please create a short title name for this image in {MAX_CHARS} characters or less."
                        }
                    ]
                }
            ]
        )
        return response.content

    def summarize_pdf(self, pdf_path):
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text()

        response = self.anthropic_client.messages.create(
            model="claude-3-sonnet-20240229",
            max_tokens=MAX_TOKENS,
            messages=[
                {"role": "user",
                    "content": f"Please create a short title name for this PDF content in {MAX_CHARS} characters or less: {text[:1000]}"}
            ]
        )
        return response.content

    def summarize_document(self, file_path):
        try:
            text = ""
            file_extension = os.path.splitext(file_path)[1].lower()

            if file_extension in ['.txt', '.rtf']:
                with open(file_path, 'r', errors='ignore') as file:
                    text = file.read()

            elif file_extension == '.docx':
                doc = Document(file_path)
                text = '\n'.join(
                    [paragraph.text for paragraph in doc.paragraphs])

            elif file_extension == '.pdf':
                with open(file_path, 'rb') as file:
                    pdf_reader = PdfReader(file)
                    text = '\n'.join([page.extract_text()
                                     for page in pdf_reader.pages])

            elif file_extension == '.pptx':
                prs = Presentation(file_path)
                text = '\n'.join(
                    [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, 'text')])

            elif file_extension in ['.xlsx', '.xls']:
                wb = load_workbook(file_path, read_only=True, data_only=True)
                text = '\n'.join([cell.value for sheet in wb.worksheets for row in sheet.iter_rows(
                ) for cell in row if cell.value])

            else:
                raise ValueError(f"Unsupported file type: {file_extension}")

            # Truncate text if it's too long
            max_chars = 1000
            truncated_text = text[:max_chars] if len(
                text) > max_chars else text

            # Use Claude API to summarize the text
            response = self.anthropic_client.messages.create(
                model="claude-3-sonnet-20240229",
                max_tokens=MAX_TOKENS,
                messages=[
                    {"role": "user", "content": f"Please create a short title name for this document content in {MAX_CHARS} characters or less: {truncated_text}"}
                ]
            )

            # Extract and return the summary
            summary = str(response.content).strip()
            # Ensure it's no longer than MAX_CHARS characters
            return summary[:MAX_CHARS]

        except Exception as e:
            logging.info(f"Error summarizing document {file_path}: {e}")
            return os.path.basename(file_path)[:MAX_CHARS]

    def get_video_metadata(self, video_path):
        try:
            file = MutagenFile(video_path)
            title = file.get('title', [''])[0]
            if title:
                return title[:MAX_CHARS]
            else:
                # If no title is found, use the filename without extension
                base_name = os.path.basename(video_path)
                name_without_ext = os.path.splitext(base_name)[0]
                return name_without_ext[:MAX_CHARS]
        except Exception as e:
            logging.info(f"Error reading video metadata: {e}")
            return os.path.basename(video_path)[:MAX_CHARS]


if __name__ == '__main__':
    event_handler = FileHandler()
    observer = Observer()
    for directory in WATCH_DIRECTORIES:
        observer.schedule(event_handler, directory, recursive=True)
    observer.start()

    try:
        while True:
            time.sleep(5)
    except KeyboardInterrupt:
        observer.stop()
    observer.join()
else:
    logging.info(f'Not running observer as running from {__name__}')
